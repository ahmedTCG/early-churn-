from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    c = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(6))
    tf2 = c.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(content):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(4)

add_slide("Project Overview", [
    "Project: Customer Churn Prediction Model",
    "",
    "Objective: Build a machine learning model to identify customers likely to become",
    "inactive in the next 30 days, enabling proactive retention efforts.",
    "",
    "Business Context:",
    "  - Customer retention is more cost-effective than acquisition",
    "  - Early identification allows targeted intervention",
    "  - Data-driven prioritization of retention resources",
    "",
    "Current Status: Production ready",
    "",
    "Key Metrics:",
    "  ROC-AUC:   90.5% (Model correctly ranks customers 90.5% of the time)",
    "  PR-AUC:    88.8% (Strong performance on imbalanced data)",
    "  F1 Score:  85.5% (Good balance of precision and recall)",
])

add_slide("Churn Definition", [
    "Definition: A customer is 'churned' if they have zero engagement for 30 days.",
    "",
    "Rationale for 30-day window:",
    "  - Balances early detection with prediction reliability",
    "  - Aligns with monthly business cycles",
    "  - Provides sufficient time for intervention",
    "",
    "Activity Events Included (Positive Engagement):",
    "  emarsys_open       Customer opened marketing email      11.4M events",
    "  emarsys_click      Customer clicked link in email       427K events",
    "  order              Customer completed a purchase        178K events",
    "  emarsys_sessions_* Customer visited website             8.4M events",
    "",
    "Activity Events Excluded (Non-Engagement):",
    "  emarsys_cancel, emarsys_hard_bounce, emarsys_soft_bounce, emarsys_unsub",
    "",
    "Note: Customers with ONLY excluded events are labeled 'weak_signal'",
])

add_slide("Data Source & Extraction", [
    "Source Database: Amazon Redshift",
    "Table: poc_dw.customer_interactions_fact",
    "",
    "Extraction Query (see data/extract_customers.sql):",
    "  - Random sample of 200,000 distinct customers",
    "  - Last 36 months of data",
    "  - Filter: incoming_outgoing = 'incoming'",
    "",
    "SQL:",
    "  WITH sampled_customers AS (",
    "      SELECT external_customerkey FROM (",
    "          SELECT DISTINCT external_customerkey",
    "          FROM poc_dw.customer_interactions_fact",
    "          WHERE event_time >= DATEADD(month, -36, CURRENT_DATE)",
    "            AND incoming_outgoing = 'incoming'",
    "      ) c ORDER BY RANDOM() LIMIT 200000",
    "  )",
    "  SELECT external_customerkey, event_time, interaction_type ...",
])

add_slide("Data Summary", [
    "Dataset Statistics:",
    "  Total Customers:           200,000",
    "  Total Events:              21,671,807",
    "  Activity Events (filtered): 20,417,637",
    "  Date Range:                Feb 2023 - Jan 2026",
    "",
    "Event Type Distribution:",
    "  emarsys_open                      11,433,845  (56.0%)",
    "  emarsys_sessions_content_url       3,505,230  (17.2%)",
    "  emarsys_sessions_content_category  3,498,651  (17.1%)",
    "  emarsys_sessions_content_tag         643,060  (3.1%)",
    "  emarsys_click                        426,756  (2.1%)",
    "  emarsys_sessions_view                318,715  (1.6%)",
    "  emarsys_sessions_purchase            280,125  (1.4%)",
    "  order                                177,635  (0.9%)",
    "",
    "Customer Segmentation:",
    "  With activity events:     187,923 (94%)",
    "  Only bounce/cancel:        12,077 (6%)",
])

add_slide("Pipeline Architecture", [
    "Pipeline: Raw CSV -> Clean -> Filter -> Partition -> Train -> Score -> Output",
    "",
    "Step  Script                  Output                      Runtime",
    "1     clean.py                cleaned_events.parquet      ~60 sec",
    "2     activity.py             activity_events.parquet     ~30 sec",
    "3     partition_activity.py   activity_events_ds/         ~60 sec",
    "4     train_lgbm.py           churn_model_lgbm.joblib     ~120 sec",
    "5     score.py                customer_scores.parquet     ~60 sec",
    "",
    "Total Runtime: ~5 minutes",
    "Peak Memory: ~3 GB",
    "",
    "Execution Command:",
    "  python scripts/run_pipeline.py",
])

add_slide("Feature Engineering", [
    "Total Features: 53",
    "",
    "1. Recency Features:",
    "   recency_days, inactive_14d, inactive_30d",
    "",
    "2. Frequency Features (30d, 60d, 90d windows):",
    "   n_events_last_Xd, active_days_last_Xd",
    "",
    "3. Event Type Counts (30d, 60d, 90d windows):",
    "   cnt_emarsys_open_last_Xd, cnt_emarsys_click_last_Xd,",
    "   cnt_order_last_Xd, cnt_emarsys_sessions_*_last_Xd",
    "",
    "4. Trend Features:",
    "   freq_trend_30_60 = n_events_30d / (n_events_60d + 1)",
    "   freq_trend_60_90 = n_events_60d / (n_events_90d + 1)",
    "",
    "5. Ratio Features:",
    "   active_ratio_30d, session_share_30d, email_share_30d,",
    "   click_to_open_30d, events_per_active_day_30d",
])

add_slide("Model Architecture", [
    "Algorithm: LightGBM (Light Gradient Boosting Machine)",
    "",
    "Why LightGBM:",
    "  - Handles large datasets efficiently",
    "  - Good performance on tabular data",
    "  - Fast training and inference",
    "",
    "Key Hyperparameters:",
    "  objective:        binary",
    "  num_leaves:       31",
    "  max_depth:        6",
    "  learning_rate:    0.05",
    "  n_estimators:     500",
    "  class_weight:     balanced",
    "  early_stopping:   50 rounds",
    "",
    "Training Strategy:",
    "  - 32 rolling monthly snapshots",
    "  - 1,946,215 training rows",
    "  - Test snapshot: 2026-01-01",
])

add_slide("Feature Importance", [
    "Top 15 Features by Importance:",
    "",
    "Rank  Feature                        Importance  Category",
    "1     recency_days                   682         Recency",
    "2     active_days_last_90d           494         Frequency",
    "3     cnt_emarsys_open_last_90d      327         Event Count",
    "4     freq_trend_60_90               324         Trend",
    "5     active_days_last_30d           194         Frequency",
    "6     cnt_emarsys_click_last_90d     178         Event Count",
    "7     active_days_last_60d           176         Frequency",
    "8     click_to_open_30d              133         Ratio",
    "9     cnt_order_last_90d             131         Event Count",
    "10    freq_trend_30_60               124         Trend",
    "",
    "Key Observations:",
    "  - Recency dominates: Days since last activity is strongest predictor",
    "  - 90-day window most important: Longer-term patterns matter",
    "  - Trends are predictive: Changes in behavior signal churn risk",
])

add_slide("Model Performance - Test Set", [
    "Test Set Configuration:",
    "  Snapshot date: 2026-01-01",
    "  Customers: 73,007",
    "  Churn rate: 52.7%",
    "",
    "Classification Metrics (threshold = 0.5):",
    "  Accuracy:   85.4%",
    "  Precision:  87.2%",
    "  Recall:     83.1%",
    "  F1 Score:   85.1%",
    "",
    "Ranking Metrics:",
    "  ROC-AUC:    90.5%  (Excellent discrimination)",
    "  PR-AUC:     88.8%  (Strong on imbalanced data)",
    "",
    "Optimal threshold for F1: 0.47 (F1 = 0.855)",
])

add_slide("Model Performance - Rolling Evaluation", [
    "12-Month Rolling Evaluation:",
    "",
    "Snapshot      Customers  Churn   ROC-AUC  PR-AUC  F1-opt",
    "2026-01-01    73,007     52.7%   0.905    0.888   0.855",
    "2025-12-02    65,708     39.7%   0.890    0.807   0.780",
    "2025-11-02    60,397     37.9%   0.886    0.784   0.769",
    "2025-10-03    58,233     37.5%   0.893    0.790   0.779",
    "2025-09-03    57,824     38.8%   0.898    0.805   0.790",
    "2025-08-04    58,120     40.5%   0.901    0.823   0.800",
    "2025-07-05    59,657     42.2%   0.899    0.831   0.810",
    "2025-06-05    60,640     43.1%   0.908    0.851   0.820",
    "",
    "Summary: ROC-AUC range 0.886-0.908 (stable, low variance)",
])

add_slide("Scoring Results", [
    "Scoring Configuration:",
    "  Snapshot: 2026-01-31",
    "  Total customers: 200,000",
    "",
    "Bucket Thresholds:",
    "  very_high:   >= 0.8  (Almost certainly will churn)",
    "  high:        >= 0.6  (Likely to churn)",
    "  medium:      >= 0.4  (At some risk)",
    "  low:         < 0.4   (Unlikely to churn)",
    "  weak_signal: N/A     (No activity data)",
    "",
    "Distribution:",
    "  very_high:    149,920  (75.0%)  Priority 1",
    "  high:          14,057  (7.0%)   Priority 2",
    "  medium:         7,721  (3.9%)   Priority 3",
    "  low:           16,225  (8.1%)   Maintain",
    "  weak_signal:   12,077  (6.0%)   Data review",
    "",
    "Note: 62% of customers are dormant (>90 days inactive)",
])

add_slide("Output Schema", [
    "Primary Output: customer_scores.parquet",
    "",
    "Column                Type      Description",
    "external_customerkey  string    Unique customer identifier",
    "snapshot_time         datetime  When scoring was performed",
    "churn_probability     float64   Model probability (0-1)",
    "churn_bucket          string    Risk category",
    "churn_action          int64     Action flag (1=intervene)",
    "",
    "Artifact Files (artifacts/ directory):",
    "  churn_model_lgbm.joblib      Trained model",
    "  feature_list_lgbm.json       Feature names (53)",
    "  train_meta_lgbm.json         Training metadata",
    "  rolling_eval_lgbm.csv        12-month evaluation",
    "  feature_importance_lgbm.csv  Feature importance",
    "  bucket_thresholds.json       Threshold configuration",
])

add_slide("Configuration Reference", [
    "Configuration File: src/churn/config.py",
    "",
    "Model Settings:",
    "  CHURN_WINDOW_DAYS = 30      # Days of inactivity = churn",
    "  LOOKBACK_DAYS = 90          # Feature calculation window",
    "  ROLLING_STEP_DAYS = 30      # Training snapshot interval",
    "  MAX_POOL_SNAPSHOTS = 0      # 0 = use all available",
    "  N_EVAL_SNAPSHOTS = 12       # Rolling evaluation periods",
    "",
    "Scoring Thresholds:",
    "  THRESH_VERY_HIGH = 0.8",
    "  THRESH_HIGH = 0.6",
    "  THRESH_MEDIUM = 0.4",
    "",
    "Activity Events:",
    "  ACTIVITY_EVENTS = {'emarsys_open', 'emarsys_click', 'order'}",
    "  ACTIVITY_PREFIX = 'emarsys_sessions_'",
])

add_slide("Data Refresh Process", [
    "Step-by-Step Refresh:",
    "",
    "1. Extract Data from Redshift:",
    "   Run query in data/extract_customers.sql",
    "   Export results to CSV",
    "",
    "2. Place File in Project:",
    "   mv ~/Downloads/export.csv 3_years_churn.csv",
    "",
    "3. Run Pipeline:",
    "   python scripts/run_pipeline.py",
    "",
    "4. Verify Output:",
    "   cat artifacts/train_meta_lgbm.json",
    "",
    "Recommended Refresh Schedule:",
    "  Weekly   - High-value customer monitoring",
    "  Monthly  - Regular retention campaigns",
    "  Quarterly - Model performance review",
])

add_slide("Model Comparison: LightGBM vs SGD", [
    "Metric          SGDClassifier   LightGBM    Improvement",
    "ROC-AUC         0.891           0.905       +1.4%",
    "PR-AUC          0.870           0.888       +1.8%",
    "F1-opt          0.848           0.855       +0.7%",
    "Training time   ~3 min          ~2 min      -33%",
    "Features        49              53          +4",
    "",
    "Why LightGBM Performs Better:",
    "  - Captures non-linear relationships",
    "  - Handles feature interactions automatically",
    "  - Better with skewed distributions",
    "  - No manual feature scaling required",
    "",
    "Decision: LightGBM selected as production model",
])

add_slide("Known Limitations", [
    "Data Limitations:",
    "  - Sample of 200K customers (not full population)",
    "  - Only 'incoming' interactions included",
    "  - 36-month history may miss longer patterns",
    "",
    "Model Limitations:",
    "  - Cannot predict for customers with no activity history",
    "  - Does not incorporate demographic or transaction value data",
    "  - Fixed 30-day churn window may not suit all use cases",
    "",
    "Operational Limitations:",
    "  - Requires manual data refresh",
    "  - No real-time scoring capability",
    "  - No automated drift detection",
])

add_slide("Next Steps", [
    "Short-term (Next Sprint):",
    "  [ ] Integrate output with CRM system",
    "  [ ] Set up automated weekly refresh",
    "  [ ] Create monitoring dashboard",
    "",
    "Medium-term (Next Quarter):",
    "  [ ] Add feature drift monitoring",
    "  [ ] Implement A/B testing for retention campaigns",
    "  [ ] Add unit tests for feature engineering",
    "",
    "Long-term (Next 6 Months):",
    "  [ ] Expand to full customer population",
    "  [ ] Add customer lifetime value features",
    "  [ ] Explore real-time scoring options",
])

add_slide("Summary", [
    "What We Built:",
    "  - Production-ready churn prediction model",
    "  - Automated 5-step pipeline",
    "  - Scoring for 200K customers",
    "",
    "Key Results:",
    "  - 90.5% ROC-AUC accuracy",
    "  - Stable performance across 12 months",
    "  - Clear risk segmentation",
    "",
    "Business Impact:",
    "  - 163,977 customers identified as at-risk",
    "  - Prioritized list for retention team",
    "  - Data-driven resource allocation",
    "",
    "",
    "Questions?",
])

prs.save('churn_model_presentation.pptx')
print("Saved: churn_model_presentation.pptx")
